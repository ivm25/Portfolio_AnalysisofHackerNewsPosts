#!/usr/bin/env python
# coding: utf-8

# ### Hacker News Posts Analysis

# *Business Problem*
# ***
# **Analyse Hacker News posts to understand which is the best time to have more comments on posts.**
# 
# **Assumptions: The dataset is assumed to largely clean. The source of data is Hacker news source on the source website.**

# In[7]:


from csv import reader
opened_file = open("hacker_news.csv",encoding="utf8")
read_file = reader(opened_file)
hn = list(read_file)


# In[8]:


#Print the first six rows to understand how the data looks and if we have the correct data or not.
print(hn[0:6])


# In[9]:


#print only the headers of the file
headers=hn[0]
print(headers)


# In[10]:


# Store the data without the headers in a new variable
hn=hn[1:]


# In[11]:


print(hn[0:5])


# ### Analysis of ask_posts, show posts and other posts

# In[12]:


#creating empty lists to store ask_posts, show_posts and other_posts
ask_posts=[]
show_posts=[]
other_posts=[]


for row in hn:
    title=row[1]
    if title.lower().startswith('ask hn'):
        ask_posts.append(row)
    elif title.lower().startswith('show hn'):
        show_posts.append(row)
    else:
        other_posts.append(row)
        


# In[13]:


print(len(show_posts))


# In[14]:


print(len(ask_posts))


# In[15]:


print(ask_posts[0:4])


# In[16]:


print(show_posts[0:4])


# **Now that we have ask posts and show posts in different lists, we will analyse the average comments per posts for both and hypothesize a reason for the differences**

# In[17]:


# Computing the average comments on ask posts

total_ask_comments=0

for row in ask_posts:
    num_contents=row[4]
    #convert num_contents ino integer values
    num_contents=int(num_contents)
    total_ask_comments += num_contents
avg_ask_comments=total_ask_comments/len(ask_posts)

print(avg_ask_comments)


# In[18]:


# Computing the average comments on show posts

total_show_comments=0

for row in show_posts:
    num_contents=row[4]
    #convert num_contents ino integer values
    num_contents=int(num_contents)
    total_show_comments += num_contents
avg_show_comments=total_show_comments/len(show_posts)

print(avg_show_comments)


# ### Analysis of the above computed average comments for ask_posts and showposts

# 1. *Ask Posts perform better with higher number of average comments*, therefore, will focus on analysing the ask posts as they generate higher number of comments on an average.
# 
# ***
#   **Next Steps**:
# 
# 2. Calculate the number of **posts created each hour**, along with number of comments.
# 
# 3. Calculate **average number of comments ask posts recieved every hour** to understand what is the best time to post/comment on ask posts
# 

# In[19]:


# initiating a new, empty list to store date,time and the number of posts created during that time
result_list = []

for post in ask_posts:
    #Use the time column and number of posts column to append the empty list      
    result_list.append(
        [post[6], int(post[4])]
    )


# In[20]:


print(result_list)


# ### As the output in the list looks correct, we will start to segregate the list by timing of the comments and create sepearate dictionaries for counts by hour and comments by hour 

# In[21]:


import datetime as dt
counts_by_hour={}

comments_by_hour={}

# store the right date format in a variable
date_format= "%m/%d/%Y %H:%M"
for each_row in result_list:
    #extract the date     
    date=each_row[0]
    #extract the coment
    comment=each_row[1]
    #to extract the hour from date variabe, use dt.strptime to parse the date and then use strftime to convert the datetime object back to a string  
    hour=dt.datetime.strptime(date,date_format).strftime("%H")
    
    if hour not in counts_by_hour:
        counts_by_hour[hour] = 1
        comments_by_hour[hour] = comment
    else:
        counts_by_hour[hour] += 1
        comments_by_hour[hour] += comment
        
 
print(counts_by_hour)


# In[22]:


print(comments_by_hour)


# **Creating a list of lists to calculate the average number of comments per post created during each hour of the day.**

# In[23]:


average_comms_per_hour=[]

for comms in comments_by_hour:
    #Divide the comments by hours by the counts by hour to have average values     
    average_comms_per_hour.append([comms, comments_by_hour[comms]/counts_by_hour[comms]])
    #store the list values in a new variable
    avg_by_hour=average_comms_per_hour


# In[24]:


# perform a sanity check on the new variable "avg_by_hour"
avg_by_hour


# **Sort the above list by average comments and then by the hour.Print the five top values to make it easier to read**

# In[25]:


swap_avg_by_hour=[]

for row in avg_by_hour:
    swap_avg_by_hour.append([row[1],row[0]])
print(swap_avg_by_hour)


# In[26]:


# Sort the list by the highest number of comments by hour

sorted_swap=sorted(swap_avg_by_hour,reverse=True)


# In[27]:


print(sorted_swap)


# *Printing top 5 hours to ask posts comments*

# In[28]:


print(sorted_swap[0:5])


# **Print out the summary of sorted list of average comments/hour/post**

# In[29]:


summary = []
for row in sorted_swap[0:5]:
    time=dt.datetime.strptime(row[1],"%H").strftime("%H:00")
    print(time)
    average=row[0]
    #create a variable using (.format()) to save the variable    
    output="{h}:{a:.2f} average comments per post".format(h=time,a=average)
    print(output)
    summary.append(output)
    


# In[30]:


type(output)


# ### Conclusion
# 
# *The above analysis demonstrates that posts around 15:00 hours on an average have a higher chance of recieving comments.*

# In[31]:


from pptx import Presentation
from pptx.util import Inches,Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# In[44]:


def create_summary(lst):
    prs=Presentation()
    im_slide=prs.slides.add_slide(prs.slide_layouts[0])
    title_slide_layout=prs.slide_layouts[0]
    slide=im_slide
                                  
    placeholder=0
    left=1
    top=1
    width=1
    height=0.5
    for row in lst:
        if placeholder==16:
            slide=prs.slides.add_slide(prs.slide_layouts[0])
            placeholder = 0
            left = 1
            top = 2.5
        if placeholder == 8 or placeholder == 16:
            left = left + 4.5
            top = 2.5
        
                                  
        textBox = slide.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(height))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        tf = shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = row
        run.font.name = "Arial"
        run.font.size = Pt(7)
        
                                  
        placeholder = placeholder + 1
                                  
        top = top + 0.5
    prs.save('Summary.pptx')


# In[45]:


create_summary(summary)

